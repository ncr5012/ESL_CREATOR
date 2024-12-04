from openai import OpenAI
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json
import os
import random
import streamlit as st
import tempfile

openai_api_key = st.secrets["openai_api_key"]
stable_diffusion_api_key = st.secrets["stable_diffusion_api_key"]

client = OpenAI(api_key=openai_api_key)

model_name = "o1-mini"


def get_target_words():
    target_words = [input(f"Word {i+1}: ").strip() for i in range(2)]
    return target_words

def generate_confusion_words(client, model_name, word):
    messages = [
        {
            "role": "user",
            "content": (
                f"You are to create an exercise for level 0 English as a Second Language learners who have absolutely no knowledge of English. "
                f"Your task is to provide two words that could be verbally or visually confused with '{word}'. "
                "The words should be easily explainable in classroom by pointing at an object or miming. "
                "The words should be single words, not phrases. Respond exclusively with a JSON Object with no additional text or labels: JSON object {\"confusion_words\": [\"word_1\", \"word_2\"]}. "
                "An example is user_input: electricity correct response: {\"confusion_words\": [\"elasticity\", \"electrician\"]}."
            )
        }
    ]
    
    # Define the JSON schema for the expected output
    json_schema = {
        "name": "confusion_words_schema",
        "schema": {
            "type": "object",
            "properties": {
                "confusion_words": {
                    "type": "array",
                    "items": {
                        "type": "string"
                    },
                    "minItems": 2,
                    "maxItems": 2
                }
            },
            "required": ["confusion_words"],
            "additionalProperties": False
        }
    }
    
    # Get response from OpenAI using the new API with response_format
    completion = client.chat.completions.create(
        model=model_name,
        messages=messages,
    )

    # Extract the assistant's reply
    confusion_words_str = completion.choices[0].message.content

    #st.write(f"Assistant response: {confusion_words_str}")
    
    # Print the assistant's response for debugging
    print(confusion_words_str)
    
    # Remove code block formatting if present
    confusion_words_str = confusion_words_str.strip()
    if confusion_words_str.startswith('```') and confusion_words_str.endswith('```'):
        # Remove the triple backticks
        confusion_words_str = confusion_words_str[3:-3].strip()
        # Remove any language identifier (e.g., 'json')
        if confusion_words_str.startswith('json'):
            confusion_words_str = confusion_words_str[4:].strip()

    # Parse the JSON string into a Python dictionary
    try:
        confusion_words_dict = json.loads(confusion_words_str)
        confusion_word_list = confusion_words_dict.get("confusion_words", [])
    except json.JSONDecodeError as e:
        print(f"JSON parsing error: {e}")
        confusion_word_list = []

    return confusion_word_list[:2]
def generate_image(word, filename):
    response = requests.post(
        f"https://api.stability.ai/v2beta/stable-image/generate/ultra",
        headers={
            "authorization": f"Bearer {stable_diffusion_api_key}",
            "accept": "image/*"
        },
        files={"none": ''},
        data={
            "prompt": f"realistic easily identifiable example of {word}",
            "output_format": "png",
        },
    )

    if response.status_code == 200:
        with open(filename, 'wb') as file:
            file.write(response.content)
    else:
        raise Exception(f"API Error: {response.status_code} - {response.text}")

def add_slide(prs, image_path, choices):
    import random

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add image on the left
    left = Inches(1)
    top = Inches(1.9)
    pic = slide.shapes.add_picture(image_path, left, top, width=Inches(4))

    # Define colors
    colors = [
        RGBColor(255, 192, 0),   # orange
        RGBColor(0, 176, 240),   # blue
        RGBColor(146, 208, 80),  # green
        RGBColor(255, 0, 0),     # red
        RGBColor(112, 48, 160),  # purple
        RGBColor(0, 112, 192),   # dark blue
        RGBColor(0, 176, 80),    # dark green
        RGBColor(255, 192, 0),   # gold
        RGBColor(255, 0, 255),   # magenta
        RGBColor(0, 176, 240)    # cyan
    ]

    shape_width = Inches(3)
    shape_height = Inches(1)
    spacing = Inches(0.5)  # Vertical spacing between choices

    num_choices = len(choices)
    slide_height = prs.slide_height
    slide_width = prs.slide_width

    # Calculate total height needed for all choices and spacings
    total_height = num_choices * shape_height + (num_choices - 1) * spacing

    # Starting y position to center choices vertically
    start_y = (slide_height - total_height) / 2

    # x position for choices, roughly 2/3 of the way to the right
    x = slide_width * 0.6

    # Shuffle colors
    random.shuffle(colors)

    for i, choice in enumerate(choices):
        # Randomly select a color
        color = colors.pop()

        # Calculate y position for each choice
        y = start_y + i * (shape_height + spacing)

        # Create rounded rectangle for each choice
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left=x,
            top=y,
            width=shape_width,
            height=shape_height
        )

        # Set fill color for the shape
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = color

        # Add text to the shape
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{chr(65 + i)}. {choice}"
        font = run.font
        font.size = Pt(24)
        font.bold = True
        font.color.rgb = RGBColor(255, 255, 255)  # White text

        # Remove the outline of the shape
        shape.line.color.rgb = color

def add_audio_to_presentation(prs, audio_path):

    # Add audio to the first slide
    if len(prs.slides) == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
    else:
        slide = prs.slides[0]

    # Add the audio file to the slide
    left = top = Inches(0)
    width = height = Inches(0.5)  # Small icon

    # Use add_movie method to add the audio
    # Note: mime_type for mp3 is 'audio/mpeg'
    shape = slide.shapes.add_movie(audio_path, left, top, width, height, mime_type='audio/mpeg')

    # Hide the audio icon during presentation (optional)
    # shape.element.getparent().remove(shape.element)

    # Note: Due to limitations of python-pptx, setting the audio to loop and play across slides
    # requires manual adjustment in PowerPoint after the presentation is generated.


def main():
    st.title("ESL sound-alike exercise creator")

    st.write("This tool allows you to create a powerpoint presentation to test a student's ability to distinguish between similar sounding/looking english words")
    st.write("Each slide consists of a picture of the word the students are to identify, along with two decoy words that sound similar (IE: Electricity, Elasticity, Electrician")
    st.write("During class, distribute A, B, and C cards, show a slide, and ask the students to hold up the letter of the word that matches the picture")
    st.write("Select the number of slides you want, then type a word in for each slide and wait 2-3 minutes for it to generate")
    st.write("FOR TESTING DO NOT SELECT more than 2x slides")
    
    number_of_slides = int(st.number_input("Number of slides", min_value=1, max_value=10, value=2, step=1))

    target_words = []
    for i in range(number_of_slides):
        word = st.text_input(f"Word {i+1}", key=f"word_{i}")
        target_words.append(word.strip())


    if st.button("Generate ESL Exercise"):
    # Define the list of audio tracks
    #audio_tracks = ['audio1.mp3', 'audio2.mp3', 'audio3.mp3']

    # Randomly select one audio track
    #selected_audio = random.choice(audio_tracks)

    # Check if the audio file exists
    #if not os.path.exists(selected_audio):
        #print(f"Audio file '{selected_audio}' not found. Please provide the audio files in the current directory.")
        #return


    # Check if all words have been entered
        if any(word == "" for word in target_words):
            st.warning("Please enter all words before generating the exercise.")
            st.stop()

            # Inform the user and show a spinner
        with st.spinner('Generating presentation... This may take 2-3 minutes.'):
            prs = Presentation()
            for i, word in enumerate(target_words):
                confusion_words = generate_confusion_words(client, model_name, word)
                all_choices = [word] + confusion_words
                random.shuffle(all_choices)  # Shuffle choices

                image_filename = f"image_{i+1}.png"
                generate_image(word, image_filename)
                add_slide(prs, image_filename, all_choices)
                st.write(f"Created slide {i}")
        #images_created = []  # List to keep track of created images

        # Add the selected audio to the presentation
        #add_audio_to_presentation(prs, selected_audio)


            # Save the presentation to a temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_pptx_file:
                prs.save(tmp_pptx_file.name)
                tmp_pptx_file_path = tmp_pptx_file.name

            # Read the saved PowerPoint file
            with open(tmp_pptx_file_path, 'rb') as f:
                pptx_data = f.read()

            # Provide a download link
            st.download_button(
                label="Download ESL Exercise PowerPoint",
                data=pptx_data,
                file_name="ESL_Exercise.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

if __name__ == "__main__":
    import random
    main()
